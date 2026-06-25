"""Fixed-template PPTX placeholder filling.

The main idea is deliberately small:

- The PPTX template owns layout, visual style, and fixed wording.
- The caller supplies values for placeholders such as ``{{content1}}``.
- This class replaces text placeholders and can turn table/chart placeholders
  into native PowerPoint tables or charts.

Supported placeholder locations:

- Text inside a shape, for example ``{{content1}}``.
- Shape names, for example a rectangle named ``{{chart1}}`` or ``chart1``.
- Shape alt text/description, for example ``{{table1}}`` or ``table1``.
"""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
import re
from typing import Any, Iterable, Literal

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Pt
except ModuleNotFoundError as exc:  # pragma: no cover - dependency guard
    raise ModuleNotFoundError(
        "python-pptx is required for ppt_template filling. "
        "Install egis-agents with the pptmaster extra, for example: "
        "`uv sync --extra pptmaster`."
    ) from exc


PLACEHOLDER_RE = re.compile(r"\{\{\s*([A-Za-z_][A-Za-z0-9_]*)\s*\}\}")


@dataclass(slots=True)
class TableFill:
    """Native PowerPoint table fill."""

    columns: list[str]
    rows: list[list[Any]]
    title: str | None = None


@dataclass(slots=True)
class ChartFill:
    """Native PowerPoint bar/line chart fill."""

    chart_type: Literal["bar", "line"]
    categories: list[str]
    series: dict[str, list[float]]
    title: str | None = None


@dataclass(slots=True)
class PlaceholderMatch:
    slide_index: int
    slide_number: int
    name: str
    source: Literal["text", "shape_name", "alt_text", "table_cell", "chart_title"]
    shape_name: str


@dataclass(slots=True)
class FillResult:
    output_path: str
    placeholders_found: list[str] = field(default_factory=list)
    filled: list[str] = field(default_factory=list)
    missing: list[str] = field(default_factory=list)
    unused_values: list[str] = field(default_factory=list)
    deleted_slide_numbers: list[int] = field(default_factory=list)
    warnings: list[str] = field(default_factory=list)


class PptTemplateFiller:
    """Fill a fixed PPTX template with text, table, and chart values."""

    def __init__(self, template_path: str | Path):
        self.template_path = Path(template_path).expanduser().resolve()
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template PPTX does not exist: {self.template_path}")
        if self.template_path.suffix.lower() != ".pptx":
            raise ValueError(f"Expected a .pptx file, got: {self.template_path}")

    def inspect_placeholders(self) -> list[PlaceholderMatch]:
        """Return placeholders found in the template without modifying it."""

        prs = Presentation(str(self.template_path))
        return list(self._iter_placeholders(prs))

    def render(
        self,
        output_path: str | Path,
        fills: dict[str, Any],
        *,
        exclude_slide_numbers: Iterable[int] | None = None,
        exclude_first: bool = False,
        exclude_last: bool = False,
        keep_unknown_placeholders: bool = True,
    ) -> FillResult:
        """Render a filled PPTX.

        Args:
            output_path: Where to save the generated PPTX.
            fills: Mapping from placeholder name to value. Values may be plain
                strings, ``TableFill``/``ChartFill`` instances, or equivalent
                dictionaries with ``type`` set to ``table`` or ``chart``.
            exclude_slide_numbers: 1-based slide numbers to remove before fill.
            exclude_first: Remove the first slide.
            exclude_last: Remove the last slide.
            keep_unknown_placeholders: If false, unresolved text placeholders
                are replaced with an empty string.
        """

        prs = Presentation(str(self.template_path))
        normalized_fills = {key: self._normalize_fill(value) for key, value in fills.items()}

        deleted = self._delete_requested_slides(
            prs,
            exclude_slide_numbers=exclude_slide_numbers,
            exclude_first=exclude_first,
            exclude_last=exclude_last,
        )

        found_names: list[str] = []
        filled: list[str] = []
        warnings: list[str] = []

        allowed_names = set(normalized_fills)

        for slide in prs.slides:
            for shape in list(slide.shapes):
                shape_placeholders = self._placeholders_for_shape(shape, allowed_names=allowed_names)
                for name, source in shape_placeholders:
                    if name not in found_names:
                        found_names.append(name)
                    if name not in normalized_fills:
                        continue

                    value = normalized_fills[name]
                    if isinstance(value, TableFill):
                        if getattr(shape, "has_table", False):
                            self._fill_existing_table_shape(shape, value)
                        else:
                            self._replace_shape_with_table(slide, shape, value)
                    elif isinstance(value, ChartFill):
                        if getattr(shape, "has_chart", False):
                            self._fill_existing_chart_shape(shape, value)
                        else:
                            self._replace_shape_with_chart(slide, shape, value)
                    else:
                        if source == "text" and hasattr(shape, "text"):
                            self._replace_text_placeholder(shape, name, str(value), keep_unknown_placeholders)
                        elif hasattr(shape, "text"):
                            self._set_shape_text(shape, str(value))
                        else:
                            warnings.append(
                                f"Placeholder {name!r} is on a non-text shape; text value was ignored."
                            )
                            continue

                    if name not in filled:
                        filled.append(name)

        if not keep_unknown_placeholders:
            self._clear_unknown_text_placeholders(prs, set(normalized_fills))

        output = Path(output_path).expanduser().resolve()
        output.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output))

        missing = [name for name in found_names if name not in normalized_fills]
        unused = [name for name in normalized_fills if name not in found_names]
        if unused:
            warnings.append(
                "These fill values were supplied but no matching placeholder was found: "
                + ", ".join(unused)
            )

        return FillResult(
            output_path=str(output),
            placeholders_found=found_names,
            filled=filled,
            missing=missing,
            unused_values=unused,
            deleted_slide_numbers=deleted,
            warnings=warnings,
        )

    def _iter_placeholders(self, prs: Any) -> Iterable[PlaceholderMatch]:
        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                for name, source in self._placeholders_for_shape(shape):
                    yield PlaceholderMatch(
                        slide_index=slide_index,
                        slide_number=slide_index + 1,
                        name=name,
                        source=source,
                        shape_name=getattr(shape, "name", ""),
                    )

    def _placeholders_for_shape(
        self,
        shape: Any,
        *,
        allowed_names: set[str] | None = None,
    ) -> list[tuple[str, Literal["text", "shape_name", "alt_text", "table_cell", "chart_title"]]]:
        matches: list[tuple[str, Literal["text", "shape_name", "alt_text", "table_cell", "chart_title"]]] = []

        if hasattr(shape, "text") and shape.text:
            for name in PLACEHOLDER_RE.findall(shape.text):
                matches.append((name, "text"))

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    for name in PLACEHOLDER_RE.findall(cell.text or ""):
                        matches.append((name, "table_cell"))

        if getattr(shape, "has_chart", False):
            chart_title = self._chart_title_text(shape.chart)
            for name in PLACEHOLDER_RE.findall(chart_title):
                matches.append((name, "chart_title"))

        shape_name = getattr(shape, "name", "") or ""
        for name in self._names_from_placeholderish_value(shape_name, allowed_names=allowed_names):
            matches.append((name, "shape_name"))

        alt_text = self._shape_alt_text(shape)
        for name in self._names_from_placeholderish_value(alt_text, allowed_names=allowed_names):
            matches.append((name, "alt_text"))

        deduped: list[tuple[str, Literal["text", "shape_name", "alt_text", "table_cell", "chart_title"]]] = []
        seen: set[tuple[str, str]] = set()
        for item in matches:
            if item not in seen:
                deduped.append(item)
                seen.add(item)
        return deduped

    def _chart_title_text(self, chart: Any) -> str:
        try:
            if not chart.has_title:
                return ""
            return chart.chart_title.text_frame.text or ""
        except (AttributeError, ValueError):
            return ""

    def _names_from_placeholderish_value(
        self,
        value: str,
        *,
        allowed_names: set[str] | None = None,
    ) -> list[str]:
        value = (value or "").strip()
        if not value:
            return []
        names = PLACEHOLDER_RE.findall(value)
        if names:
            return names
        for prefix in ("placeholder:", "placeholder=", "ph:", "ph="):
            if value.lower().startswith(prefix):
                candidate = value[len(prefix):].strip()
                if re.fullmatch(r"[A-Za-z_][A-Za-z0-9_]*", candidate):
                    return [candidate]
        if allowed_names and value in allowed_names:
            return [value]
        return []

    def _shape_alt_text(self, shape: Any) -> str:
        try:
            c_nv_pr = shape._element.nvSpPr.cNvPr  # noqa: SLF001 - python-pptx has no public alt-text API.
        except AttributeError:
            try:
                c_nv_pr = shape._element.nvPicPr.cNvPr  # noqa: SLF001
            except AttributeError:
                return ""
        return c_nv_pr.get("descr") or ""

    def _normalize_fill(self, value: Any) -> str | TableFill | ChartFill:
        if isinstance(value, (TableFill, ChartFill, str)):
            return value

        if isinstance(value, dict):
            fill_type = str(value.get("type", "")).lower()
            if fill_type == "table":
                return TableFill(
                    columns=[str(col) for col in value.get("columns", [])],
                    rows=[[cell for cell in row] for row in value.get("rows", [])],
                    title=value.get("title"),
                )
            if fill_type == "chart":
                raw_series = value.get("series", {})
                if isinstance(raw_series, list):
                    raw_series = {
                        str(item.get("name", f"series_{idx + 1}")): item.get("values", [])
                        for idx, item in enumerate(raw_series)
                        if isinstance(item, dict)
                    }
                return ChartFill(
                    chart_type=str(value.get("chart_type", "bar")).lower(),  # type: ignore[arg-type]
                    categories=[str(category) for category in value.get("categories", [])],
                    series={
                        str(name): [float(v) for v in values]
                        for name, values in dict(raw_series).items()
                    },
                    title=value.get("title"),
                )
            if "text" in value:
                return str(value["text"])

        return str(value)

    def _delete_requested_slides(
        self,
        prs: Any,
        *,
        exclude_slide_numbers: Iterable[int] | None,
        exclude_first: bool,
        exclude_last: bool,
    ) -> list[int]:
        slide_count = len(prs.slides)
        to_delete = set(int(n) for n in (exclude_slide_numbers or []) if int(n) > 0)
        if exclude_first and slide_count:
            to_delete.add(1)
        if exclude_last and slide_count:
            to_delete.add(slide_count)

        for slide_number in sorted(to_delete, reverse=True):
            if slide_number < 1 or slide_number > len(prs.slides):
                continue
            slide_id = prs.slides._sldIdLst[slide_number - 1]  # noqa: SLF001 - no public delete API.
            prs.slides._sldIdLst.remove(slide_id)  # noqa: SLF001

        return sorted(n for n in to_delete if 1 <= n <= slide_count)

    def _replace_text_placeholder(
        self,
        shape: Any,
        name: str,
        value: str,
        keep_unknown_placeholders: bool,
    ) -> None:
        marker_re = re.compile(r"\{\{\s*" + re.escape(name) + r"\s*\}\}")
        replaced_in_run = False
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if marker_re.search(run.text):
                    run.text = marker_re.sub(value, run.text)
                    replaced_in_run = True

        if replaced_in_run:
            if not keep_unknown_placeholders:
                self._clear_unknown_placeholders_in_text_frame(shape.text_frame, {name})
            return

        new_text = marker_re.sub(value, shape.text)
        if not keep_unknown_placeholders:
            new_text = PLACEHOLDER_RE.sub("", new_text)
        self._set_shape_text(shape, new_text)

    def _set_shape_text(self, shape: Any, value: str) -> None:
        shape.text = value
        if not getattr(shape, "text_frame", None):
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = run.font.size or Pt(14)

    def _clear_unknown_text_placeholders(self, prs: Any, known_names: set[str]) -> None:
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text") or "{{" not in shape.text:
                    continue

                def replacement(match: re.Match[str]) -> str:
                    return match.group(0) if match.group(1) in known_names else ""

                self._set_shape_text(shape, PLACEHOLDER_RE.sub(replacement, shape.text))

    def _clear_unknown_placeholders_in_text_frame(self, text_frame: Any, known_names: set[str]) -> None:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = PLACEHOLDER_RE.sub(
                    lambda match: match.group(0) if match.group(1) in known_names else "",
                    run.text,
                )

    def _fill_existing_table_shape(self, shape: Any, table_fill: TableFill) -> None:
        if not table_fill.columns:
            raise ValueError("TableFill.columns must not be empty")

        table = shape.table
        data_rows = len(table_fill.rows)
        include_header = len(table.rows) >= data_rows + 1
        required_rows = data_rows + 1 if include_header else data_rows
        required_cols = len(table_fill.columns)
        if len(table.rows) < required_rows or len(table.columns) < required_cols:
            raise ValueError(
                "Template table is too small for the supplied data: "
                f"needs at least {required_rows} rows x {required_cols} columns, "
                f"but template has {len(table.rows)} rows x {len(table.columns)} columns."
            )

        row_offset = 0
        if include_header:
            for col_idx, col_name in enumerate(table_fill.columns):
                self._set_cell_text_preserve_style(table.cell(0, col_idx), str(col_name))
            row_offset = 1

        for row_idx, row_values in enumerate(table_fill.rows, start=row_offset):
            for col_idx in range(required_cols):
                value = row_values[col_idx] if col_idx < len(row_values) else ""
                self._set_cell_text_preserve_style(table.cell(row_idx, col_idx), str(value))

        # Clear leftover cells inside the template's data area. Widths, heights,
        # fills, borders, and fonts stay owned by the template.
        for row_idx in range(required_rows, len(table.rows)):
            for col_idx in range(required_cols):
                self._set_cell_text_preserve_style(table.cell(row_idx, col_idx), "")

    def _set_cell_text_preserve_style(self, cell: Any, value: str) -> None:
        text_frame = cell.text_frame
        paragraph = text_frame.paragraphs[0]
        if paragraph.runs:
            paragraph.runs[0].text = value
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            cell.text = value

    def _replace_shape_with_table(self, slide: Any, shape: Any, table_fill: TableFill) -> None:
        if not table_fill.columns:
            raise ValueError("TableFill.columns must not be empty")

        rows = max(1, len(table_fill.rows)) + 1
        cols = len(table_fill.columns)
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        self._remove_shape(shape)

        graphic_frame = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = graphic_frame.table
        for col_idx, col_name in enumerate(table_fill.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            self._style_cell_text(cell, bold=True)

        for row_idx, row_values in enumerate(table_fill.rows, start=1):
            for col_idx in range(cols):
                value = row_values[col_idx] if col_idx < len(row_values) else ""
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                self._style_cell_text(cell, bold=False)

    def _fill_existing_chart_shape(self, shape: Any, chart_fill: ChartFill) -> None:
        chart_data = self._chart_data_from_fill(chart_fill)
        chart = shape.chart
        chart.replace_data(chart_data)
        if chart_fill.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_fill.title
        elif chart.has_title and PLACEHOLDER_RE.search(self._chart_title_text(chart)):
            chart.has_title = False

    def _replace_shape_with_chart(self, slide: Any, shape: Any, chart_fill: ChartFill) -> None:
        if chart_fill.chart_type not in {"bar", "line"}:
            raise ValueError(f"Unsupported chart_type: {chart_fill.chart_type}")
        chart_data = self._chart_data_from_fill(chart_fill)

        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        self._remove_shape(shape)

        chart_type = (
            XL_CHART_TYPE.COLUMN_CLUSTERED
            if chart_fill.chart_type == "bar"
            else XL_CHART_TYPE.LINE_MARKERS
        )
        frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
        chart = frame.chart
        chart.has_legend = len(chart_fill.series) > 1
        chart.has_title = bool(chart_fill.title)
        if chart_fill.title:
            chart.chart_title.text_frame.text = chart_fill.title

    def _chart_data_from_fill(self, chart_fill: ChartFill) -> CategoryChartData:
        if not chart_fill.categories:
            raise ValueError("ChartFill.categories must not be empty")
        if not chart_fill.series:
            raise ValueError("ChartFill.series must not be empty")

        category_count = len(chart_fill.categories)
        for series_name, values in chart_fill.series.items():
            if len(values) != category_count:
                raise ValueError(
                    f"Series {series_name!r} has {len(values)} values, "
                    f"but categories has {category_count} values."
                )

        chart_data = CategoryChartData()
        chart_data.categories = chart_fill.categories
        for series_name, values in chart_fill.series.items():
            chart_data.add_series(series_name, values)
        return chart_data

    def _style_cell_text(self, cell: Any, *, bold: bool) -> None:
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = bold
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(32, 37, 45)

    def _remove_shape(self, shape: Any) -> None:
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PLACEHOLDER:
            # Placeholder shapes can still be removed through the XML tree.
            pass
        parent = shape._element.getparent()  # noqa: SLF001
        parent.remove(shape._element)  # noqa: SLF001
